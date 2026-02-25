import os
from typing import List
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader
)

def load_document(file_path: str) -> List[Document]:
    """
    Loads documents based on file extension.
    Supports .pdf, .docx, .pptx, .txt, .md
    Simplified for Vercel deployment (no heavy Unstructured dependencies).
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
            return loader.load()
        elif ext in [".docx", ".doc"]:
            loader = Docx2txtLoader(file_path)
            return loader.load()
        elif ext in [".pptx", ".ppt"]:
            # Lightweight PPTX extraction using python-pptx directly
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            return [Document(page_content="\n".join(full_text), metadata={"source": file_path})]
        elif ext in [".txt", ".md"]:
            loader = TextLoader(file_path)
            return loader.load()
        else:
            raise ValueError(f"Unsupported file extension: {ext}")
            
    except Exception as e:
        print(f"Error loading document {file_path}: {e}")
        # Fallback to a simple text load if it's text-like but unknown
        try:
            loader = TextLoader(file_path)
            return loader.load()
        except:
            return []
